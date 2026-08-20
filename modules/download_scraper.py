import contextlib
import csv
import os
import re
import shutil
import traceback
from collections.abc import Callable

import fitz  # PyMuPDF (já está no requirements.txt)
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyQt6.QtCore import QThread, pyqtSignal

from modules.config import PALAVRAS_CHAVE, PERFIL_EDGE_DOWNLOAD, get_coupa_base_url
from modules.playwright_pool import PlaywrightContextManager

# O Coupa usa 2 componentes de UI diferentes para anexo, cada um com seu
# próprio seletor - não precisam de escopo por posição/container pra não se
# confundir, já que não compartilham classe/atributo:
# - Anexos da seção "Anexos" (Justificativa, lá em cima): span.attachment-file
# - Anexo ligado ao item do carrinho: span com role="button" e
#   aria-label="Anexo de arquivo de <nome>" (HTML real inspecionado no Coupa),
#   sem a classe attachment-file.
_SELETOR_ANEXO_TOPO = "span.attachment-file"
_SELETOR_ANEXO_CARRINHO = "span[aria-label^='Anexo de arquivo']"


class DownloadScraper:
    def __init__(self, requisicoes: list[str], pasta_download: str):
        # Nunca processa a mesma requisição duas vezes (ex: uma requisição com
        # 2 pedidos aparece 2x na lista importada da Aba 1) - baixar a mesma
        # página/anexo de novo é desperdício e ainda gera arquivo duplicado
        # (analisar_arquivo salva o 2º como "REQ_1.pdf").
        self.requisicoes = list(dict.fromkeys(r.strip() for r in requisicoes if r.strip()))
        self.pasta_download = pasta_download
        self.arquivos_salvos_na_execucao: list[str] = []
        self.requisicoes_sem_arquivos: list[str] = []
        self.cancelado = False
        self._log_callback: Callable[[str], None] | None = None

    @staticmethod
    def normalizar_texto(texto: str) -> str:
        texto = os.path.normpath(texto or "")
        texto = texto.replace("\\", " ")
        texto = texto.replace("/", " ")
        texto = texto.strip()
        return re.sub(r"\s+", " ", texto).strip().lower()

    def contem_de_acordo(self, texto: str) -> bool:
        texto = texto.replace("_", " ").replace("-", " ")
        return "de acordo" in self.normalizar_texto(texto)

    def _extrair_texto_log(self, msg: str) -> None:
        """Melhoria 8: loga erro via _log_callback ou logger padrão."""
        try:
            if self._log_callback:
                self._log_callback(msg)
            else:
                import logging
                logging.getLogger(__name__).warning(msg)
        except Exception:
            pass  # best-effort: nem o log de erro deve derrubar o fluxo

    def extrair_texto(self, caminho_arquivo: str) -> str:
        """Extrai texto de arquivos PDF, DOCX, TXT, CSV, XLSX, PPTX.

        Melhoria 8: logging via callback em vez de print().
        """
        extensao = os.path.splitext(caminho_arquivo)[1].lower()
        texto = ""
        try:
            if extensao == ".pdf":
                with fitz.open(caminho_arquivo) as pdf:
                    for pagina in pdf:
                        texto += pagina.get_text() or ""
            elif extensao == ".docx":
                documento = Document(caminho_arquivo)
                for paragrafo in documento.paragraphs:
                    texto += paragrafo.text + "\n"
            elif extensao == ".txt":
                with open(caminho_arquivo, encoding="utf-8", errors="ignore") as f:
                    texto = f.read()
            elif extensao == ".csv":
                with open(caminho_arquivo, newline="", encoding="utf-8", errors="ignore") as f:
                    leitor = csv.reader(f)
                    for linha in leitor:
                        texto += " ".join(map(str, linha)) + "\n"
            elif extensao == ".xlsx":
                workbook = load_workbook(caminho_arquivo, data_only=True)
                for aba in workbook.worksheets:
                    for linha in aba.iter_rows(values_only=True):
                        for celula in linha:
                            if celula is not None:
                                texto += str(celula) + " "
            elif extensao == ".pptx":
                apresentacao = Presentation(caminho_arquivo)
                for slide in apresentacao.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texto += shape.text + "\n"
        except Exception as e:
            self._extrair_texto_log(
                f"[DownloadScraper] Erro ao extrair texto de "
                f"'{os.path.basename(caminho_arquivo)}': {str(e)}"
            )
            return ""
        return texto.lower()

    def arquivo_tem_formato_invalido(self, texto: str) -> bool:
        texto_normalizado = self.normalizar_texto(texto)
        if "de acordo exceto" in texto_normalizado:
            return True

        header_terms = ["de:", "para:", "assunto:", "cc:", "interno"]
        if all(term in texto_normalizado for term in header_terms):
            return True

        return "escopos ajustados" in texto_normalizado and "de acordo" in texto_normalizado

    def analisar_arquivo(
        self,
        caminho_arquivo: str,
        req_num: str,
        nome_original: str,
        log_callback=None,
    ) -> tuple[bool, str | None]:
        extensao = os.path.splitext(caminho_arquivo)[1]
        destino = os.path.normpath(os.path.join(self.pasta_download, f"{req_num}{extensao}"))

        contador = 1
        while os.path.exists(destino):
            destino = os.path.normpath(os.path.join(self.pasta_download, f"{req_num}_{contador}{extensao}"))
            contador += 1

        if self.contem_de_acordo(nome_original):
            return False, "nome"

        texto = self.extrair_texto(caminho_arquivo)
        texto_normalizado = self.normalizar_texto(texto)

        if texto_normalizado == "":
            nome_normalizado = self.normalizar_texto(nome_original)
            if any(self.normalizar_texto(palavra) in nome_normalizado for palavra in PALAVRAS_CHAVE):
                shutil.copy(caminho_arquivo, destino)
                self.arquivos_salvos_na_execucao.append(destino)
                return True, destino
            return False, "vazio"

        if self.arquivo_tem_formato_invalido(texto_normalizado):
            return False, "formato"

        for palavra in PALAVRAS_CHAVE:
            if self.normalizar_texto(palavra) in texto_normalizado:
                shutil.copy(caminho_arquivo, destino)
                self.arquivos_salvos_na_execucao.append(destino)
                return True, destino

        if any(self.normalizar_texto(palavra) in self.normalizar_texto(nome_original) for palavra in PALAVRAS_CHAVE):
            shutil.copy(caminho_arquivo, destino)
            self.arquivos_salvos_na_execucao.append(destino)
            return True, destino

        return False, "palavra"

    async def run(self, log_callback, progress_req_callback, progress_down_callback) -> bool:
        log_callback("Iniciando Edge com Perfil Isolado...")
        user_data_dir = str(PERFIL_EDGE_DOWNLOAD)  # Item 16: path centralizado em config.py
        PERFIL_EDGE_DOWNLOAD.mkdir(parents=True, exist_ok=True)

        try:
            async with PlaywrightContextManager(user_data_dir=user_data_dir) as context:
                return await self._processar(context, log_callback, progress_req_callback, progress_down_callback)
        except Exception as e:
            log_callback(f"ERRO CRÍTICO AO INICIAR EDGE: {str(e)}")
            return False

    async def _baixar_e_validar_anexos(
        self, page, elementos, req: str, log_callback, progress_down_callback,
    ) -> int:
        """Baixa cada elemento de anexo e valida com analisar_arquivo (nome + conteúdo).

        Retorna quantos arquivos foram reconhecidos como orçamento e salvos.
        """
        total = len(elementos)
        salvos = 0
        for i, el in enumerate(elementos, 1):
            if self.cancelado:
                break
            nome = await el.inner_text()

            if self.contem_de_acordo(nome):
                progress_down_callback(int((i / total) * 100))
                continue

            arquivo_temporario = os.path.join(self.pasta_download, f"temp_{nome}")
            try:
                async with page.expect_download(timeout=30000) as download_info:
                    await el.click()
                download = await download_info.value
                await download.save_as(arquivo_temporario)

                extensoes_suportadas = (".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".txt")
                if nome.lower().endswith(extensoes_suportadas):
                    sucesso, motivo = self.analisar_arquivo(arquivo_temporario, req, nome)
                    if sucesso:
                        salvos += 1
                    else:
                        log_callback(
                            f"ℹ️ Anexo '{nome}' da requisição #{req} não reconhecido como "
                            f"orçamento (motivo: {motivo})."
                        )
                else:
                    log_callback(f"ℹ️ Anexo '{nome}' da requisição #{req} ignorado: extensão não suportada.")
            except Exception as e:
                log_callback(f"⚠️ Erro ao baixar anexo '{nome}' da requisição #{req}: {str(e)}")
            finally:
                if os.path.exists(arquivo_temporario):
                    os.remove(arquivo_temporario)

            progress_down_callback(int((i / total) * 100))
        return salvos

    async def _localizar_anexos_carrinho(self, page):
        """Clica em 'Itens do carrinho' e retorna os anexos ligados aos itens do carrinho.

        Usa um seletor próprio (span com aria-label "Anexo de arquivo de...",
        ver _SELETOR_ANEXO_CARRINHO) que não é compartilhado com a seção de
        Anexos de cima, então não precisa restringir por posição/container -
        os dois seletores já não se confundem.
        """
        aba_carrinho = await page.query_selector("a:has-text('Itens do carrinho')")
        if not aba_carrinho:
            return []
        await aba_carrinho.click()
        with contextlib.suppress(Exception):
            await page.wait_for_selector(_SELETOR_ANEXO_CARRINHO, timeout=3000)
        return await page.query_selector_all(_SELETOR_ANEXO_CARRINHO)

    async def _processar(self, context, log_callback, progress_req_callback, progress_down_callback) -> bool:
        coupa_base_url = get_coupa_base_url()
        pages = context.pages
        page = pages[0] if pages else await context.new_page()

        total_reqs = len(self.requisicoes)
        for idx, req in enumerate(self.requisicoes, 1):
            if self.cancelado:
                break

            url = f"{coupa_base_url.rstrip('/')}/requisition_headers/{req.strip()}"
            log_callback(f"📂 Processando requisição #{req}...")

            try:
                try:
                    await page.goto(url, wait_until="load", timeout=60000)
                except Exception:
                    log_callback("⏳ Página demorou para carregar (possível tela de login)...")

                if "login" in page.url.lower() or "sso" in page.url.lower():
                    log_callback("⚠️ Realize o login no Edge se necessário...")
                    await page.wait_for_url(
                        lambda u: "login" not in u.lower() and "sso" not in u.lower(),
                        timeout=300000,
                    )
                    await page.wait_for_load_state("networkidle", timeout=15000)
                    await page.goto(url, wait_until="load", timeout=60000)

                # 1. Anexo do item do carrinho primeiro - é sempre a versão mais
                # atualizada do orçamento. Só cai para a seção de Anexos (abaixo)
                # se o carrinho não tiver arquivo ou o arquivo não validar como orçamento.
                arquivos_salvos_no_req = 0
                anexos_carrinho = await self._localizar_anexos_carrinho(page)
                log_callback(f"🔎 #{req}: {len(anexos_carrinho)} anexo(s) no item do carrinho.")

                if anexos_carrinho:
                    arquivos_salvos_no_req = await self._baixar_e_validar_anexos(
                        page, anexos_carrinho, req, log_callback, progress_down_callback,
                    )
                    if arquivos_salvos_no_req > 0:
                        log_callback(f"✅ Orçamento encontrado no item do carrinho da requisição #{req}.")

                if arquivos_salvos_no_req == 0:
                    try:
                        await page.wait_for_selector(_SELETOR_ANEXO_TOPO, timeout=10000)
                    except Exception as e:
                        log_callback(f"ℹ️ Nenhum anexo encontrado na seção de Anexos para #{req}: {str(e)}")

                    anexos_topo = await page.query_selector_all(_SELETOR_ANEXO_TOPO)
                    if not anexos_topo:
                        self.requisicoes_sem_arquivos.append(req)
                        progress_req_callback(int((idx / total_reqs) * 100))
                        continue

                    arquivos_salvos_no_req = await self._baixar_e_validar_anexos(
                        page, anexos_topo, req, log_callback, progress_down_callback,
                    )

                if arquivos_salvos_no_req == 0:
                    self.requisicoes_sem_arquivos.append(req)

            except Exception as e:
                log_callback(f"❌ Erro ao processar #{req}: {str(e)}")
                log_callback(f"   Detalhes: {traceback.format_exc()}")

            progress_req_callback(int((idx / total_reqs) * 100))

        return not self.cancelado


class DownloadWorker(QThread):
    log_signal = pyqtSignal(str)
    progress_req_signal = pyqtSignal(int)
    progress_down_signal = pyqtSignal(int)
    finished_signal = pyqtSignal(bool, list, list)

    def __init__(self, requisicoes: list[str], pasta_download: str):
        super().__init__()
        self.scraper = DownloadScraper(requisicoes, pasta_download)
        # Melhoria 8: expõe o log_callback ao scraper para que extrair_texto
        # possa reportar erros via UI em vez de print().
        self.scraper._log_callback = self.log_signal.emit

    def run(self):
        """Melhoria 6: try/except garante que finished_signal SEMPRE seja emitido."""
        import asyncio
        import traceback as tb

        loop = None
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            sucesso = loop.run_until_complete(
                self.scraper.run(
                    self.log_signal.emit,
                    self.progress_req_signal.emit,
                    self.progress_down_signal.emit,
                )
            )
            self.finished_signal.emit(
                sucesso,
                self.scraper.arquivos_salvos_na_execucao,
                self.scraper.requisicoes_sem_arquivos,
            )
        except Exception as e:
            self.log_signal.emit(f"❌ ERRO CRÍTICO na thread de download: {str(e)}")
            with contextlib.suppress(Exception):
                # best-effort: não deixa o log do traceback mascarar o erro original
                self.log_signal.emit(tb.format_exc())
            # Garante que a UI não fique presa com botões desabilitados
            self.finished_signal.emit(False, [], self.scraper.requisicoes_sem_arquivos)
        finally:
            if loop is not None:
                with contextlib.suppress(Exception):
                    loop.close()  # best-effort: loop pode já estar fechado

    def cancelar(self) -> None:
        self.scraper.cancelado = True

